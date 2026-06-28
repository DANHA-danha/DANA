import os
import re
import nbformat
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from anthropic import Anthropic
from dotenv import load_dotenv
from datetime import datetime

load_dotenv()
client = Anthropic()

def extract_notebook_content(notebook_path):
    with open(notebook_path, 'r', encoding='utf-8') as f:
        nb = nbformat.read(f, as_version=4)

    content = []
    for cell in nb.cells:
        if cell.cell_type == 'code':
            if cell.source.strip():
                content.append(f"[코드]\n{cell.source}")
            for output in cell.outputs:
                if output.output_type == 'stream':
                    content.append(f"[출력]\n{output.text}")
                elif output.output_type == 'display_data':
                    if 'text/plain' in output.data:
                        content.append(f"[데이터]\n{output.data['text/plain']}")
                elif output.output_type == 'execute_result':
                    if 'text/plain' in output.data:
                        content.append(f"[결과]\n{output.data['text/plain']}")
        elif cell.cell_type == 'markdown':
            content.append(f"[분석 메모]\n{cell.source}")

    return "\n\n".join(content)


def extract_csv_content(file_path):
    import pandas as pd
    df = pd.read_csv(file_path, encoding='utf-8')
    content = []
    content.append(f"[데이터 개요]\n행: {len(df):,}건, 열: {len(df.columns)}개")
    content.append(f"[컬럼 목록]\n{', '.join(df.columns.tolist())}")

    desc = df.describe(include='all').to_string()
    content.append(f"[기술 통계]\n{desc}")

    content.append(f"[상위 10행 샘플]\n{df.head(10).to_string()}")

    for col in df.select_dtypes(include='number').columns[:10]:
        content.append(f"[수치 분포 - {col}]\n평균: {df[col].mean():.2f}, 중앙값: {df[col].median():.2f}, 최소: {df[col].min()}, 최대: {df[col].max()}")

    for col in df.select_dtypes(include='object').columns[:5]:
        vc = df[col].value_counts().head(10)
        content.append(f"[범주 분포 - {col}]\n{vc.to_string()}")

    return "\n\n".join(content), df


def extract_excel_content(file_path):
    import pandas as pd
    xls = pd.ExcelFile(file_path)
    content = []
    all_dfs = []

    for sheet in xls.sheet_names[:5]:
        df = pd.read_excel(xls, sheet_name=sheet)
        all_dfs.append(df)
        content.append(f"[시트: {sheet}]\n행: {len(df):,}건, 열: {len(df.columns)}개")
        content.append(f"[컬럼]\n{', '.join(df.columns.tolist())}")
        desc = df.describe(include='all').to_string()
        content.append(f"[기술 통계]\n{desc}")
        content.append(f"[상위 10행]\n{df.head(10).to_string()}")

        for col in df.select_dtypes(include='number').columns[:10]:
            content.append(f"[수치 - {col}]\n평균: {df[col].mean():.2f}, 중앙값: {df[col].median():.2f}")

    return "\n\n".join(content), all_dfs[0] if all_dfs else None


def extract_csv_chart_data(df):
    datasets = []
    for col in df.select_dtypes(include='object').columns[:3]:
        vc = df[col].value_counts().head(8)
        if len(vc) >= 2:
            datasets.append({
                'title': f'{col}별 분포',
                'labels': vc.index.tolist(),
                'values': vc.values.tolist(),
                'unit': '건',
            })

    for col in df.select_dtypes(include='number').columns[:3]:
        if df[col].nunique() > 1:
            for cat_col in df.select_dtypes(include='object').columns[:1]:
                grouped = df.groupby(cat_col)[col].mean().head(8)
                if len(grouped) >= 2:
                    datasets.append({
                        'title': f'{cat_col}별 평균 {col}',
                        'labels': [str(x) for x in grouped.index.tolist()],
                        'values': [round(x, 2) for x in grouped.values.tolist()],
                        'unit': '',
                    })
                break

    return datasets


def setup_korean_font():
    for name in ['AppleGothic', 'Malgun Gothic', 'NanumGothic']:
        matches = [f for f in fm.fontManager.ttflist if name in f.name]
        if matches:
            plt.rcParams['font.family'] = name
            break
    plt.rcParams['axes.unicode_minus'] = False


def extract_chart_data(notebook_path):
    with open(notebook_path, 'r', encoding='utf-8') as f:
        nb = nbformat.read(f, as_version=4)

    datasets = []
    for cell in nb.cells:
        if cell.cell_type != 'code':
            continue
        for output in cell.outputs:
            text = ''
            if output.output_type == 'stream':
                text = output.text
            elif output.output_type in ('execute_result', 'display_data'):
                text = output.data.get('text/plain', '')
            if not text.strip():
                continue

            lines = text.strip().split('\n')
            title = lines[0].strip().rstrip(':')

            count_labels, count_values = [], []
            pct_labels, pct_values = [], []

            for line in lines[1:]:
                stripped = line.strip()
                if not stripped or stripped.startswith(('-', '▶', '▼', '▲', '⚠')):
                    if not re.match(r'^[-\s]*\S+.*[:：]\s*[\d,]', stripped):
                        continue

                m_pct = re.match(r'^[-\s]*(.+?)[:：]\s*([\d,]+(?:\.\d+)?)\s*%', stripped)
                m_count = re.match(r'^[-\s]*(.+?)[:：]\s*([\d,]+(?:\.\d+)?)\s*건', stripped)

                if m_pct:
                    label = m_pct.group(1).strip()
                    if '전년' not in label and '동기' not in label:
                        pct_labels.append(label)
                        pct_values.append(float(m_pct.group(2).replace(',', '')))
                elif m_count:
                    label = m_count.group(1).strip()
                    skip_keywords = ['총', '합계', '평균', '월평균', '상반기', '하반기', '전체']
                    if not any(kw in label for kw in skip_keywords):
                        count_labels.append(label)
                        count_values.append(float(m_count.group(2).replace(',', '')))

            if len(pct_labels) >= 2:
                datasets.append({
                    'title': title,
                    'labels': pct_labels,
                    'values': pct_values,
                    'unit': '%',
                })
            if len(count_labels) >= 2:
                datasets.append({
                    'title': title,
                    'labels': count_labels,
                    'values': count_values,
                    'unit': '건',
                })
    return datasets


def _style_ax(ax):
    ax.spines[['top', 'right', 'left']].set_visible(False)
    ax.tick_params(left=False, labelsize=10)
    ax.set_axisbelow(True)
    ax.yaxis.grid(True, color='#F3F4F6', linewidth=0.8)


def generate_charts_figures(datasets):
    setup_korean_font()

    pastel = ['#93C5FD', '#60A5FA', '#3B82F6', '#2563EB', '#1D4ED8', '#1E40AF', '#1E3A8A']
    soft = ['#BFDBFE', '#93C5FD', '#60A5FA', '#3B82F6', '#2563EB', '#1D4ED8', '#1E40AF']

    figures = []
    for ds in datasets:
        n = len(ds['labels'])
        total = sum(ds['values'])
        is_pct = ds['unit'] == '%'
        is_trend = any(kw in ds['title'] for kw in ['월별', '추이', '연도', '연별'])
        is_small_pct = is_pct and n <= 5 and total < 200

        if is_small_pct:
            fig, (ax_pie, ax_bar) = plt.subplots(1, 2, figsize=(12, 5),
                                                  gridspec_kw={'width_ratios': [1, 1.3]})
            colors_pie = pastel[:n]
            wedges, texts, autotexts = ax_pie.pie(
                ds['values'], labels=ds['labels'], autopct='%1.1f%%',
                colors=colors_pie, startangle=90, pctdistance=0.75,
                wedgeprops=dict(width=0.45, edgecolor='white', linewidth=2))
            for t in autotexts:
                t.set_fontsize(10)
                t.set_fontweight('bold')
            for t in texts:
                t.set_fontsize(9)
            ax_pie.set_title('비중', fontsize=12, fontweight='bold', pad=10)

            bars = ax_bar.barh(ds['labels'][::-1], ds['values'][::-1],
                               color=colors_pie[::-1], height=0.5, edgecolor='white', linewidth=1.5)
            for bar, val in zip(bars, ds['values'][::-1]):
                ax_bar.text(bar.get_width() + max(ds['values'])*0.03,
                            bar.get_y() + bar.get_height()/2,
                            f'{val}%', va='center', fontsize=12, fontweight='bold', color='#374151')
            ax_bar.spines[['top', 'right', 'bottom']].set_visible(False)
            ax_bar.tick_params(bottom=False, labelbottom=False, left=False, labelsize=10)
            ax_bar.set_title('항목별 수치', fontsize=12, fontweight='bold', pad=10)

            fig.suptitle(ds['title'], fontsize=15, fontweight='bold', y=1.02, color='#111827')
            fig.tight_layout()
            figures.append(fig)

        elif is_trend:
            fig, ax = plt.subplots(figsize=(11, 5))
            ax.plot(ds['labels'], ds['values'], marker='o', color='#1E3A8A',
                    linewidth=3, markersize=10, markerfacecolor='white',
                    markeredgecolor='#1E3A8A', markeredgewidth=2.5, zorder=5)
            ax.fill_between(ds['labels'], ds['values'], alpha=0.08, color='#1E3A8A')

            max_val = max(ds['values'])
            max_idx = ds['values'].index(max_val)
            for j, (x, y) in enumerate(zip(ds['labels'], ds['values'])):
                color = '#1E3A8A' if j == max_idx else '#374151'
                weight = 'bold'
                ax.text(x, y + max_val * 0.03, f'{y:,.0f}',
                        ha='center', va='bottom', fontsize=11, fontweight=weight, color=color)
            if max_idx >= 0:
                ax.scatter([ds['labels'][max_idx]], [max_val], s=150,
                           color='#1E3A8A', zorder=6, alpha=0.3)

            _style_ax(ax)
            ax.spines['bottom'].set_visible(True)
            ax.spines['bottom'].set_color('#E5E7EB')
            ax.set_title(ds['title'], fontsize=15, fontweight='bold', pad=18, color='#111827')
            fig.tight_layout()
            figures.append(fig)

        elif is_pct:
            fig, ax = plt.subplots(figsize=(11, 5))
            sorted_idx = sorted(range(n), key=lambda k: ds['values'][k], reverse=True)
            labels_s = [ds['labels'][k] for k in sorted_idx]
            values_s = [ds['values'][k] for k in sorted_idx]
            colors_s = soft[:n]

            bars = ax.barh(labels_s[::-1], values_s[::-1],
                           color=colors_s[::-1], height=0.55, edgecolor='white', linewidth=1.5)
            for bar, val in zip(bars, values_s[::-1]):
                ax.text(bar.get_width() + max(values_s)*0.02,
                        bar.get_y() + bar.get_height()/2,
                        f'{val}%', va='center', fontsize=12, fontweight='bold', color='#374151')
            ax.spines[['top', 'right', 'bottom']].set_visible(False)
            ax.tick_params(bottom=False, labelbottom=False, left=False, labelsize=10)
            ax.set_title(ds['title'], fontsize=15, fontweight='bold', pad=18, color='#111827')

            top_label = labels_s[0]
            top_val = values_s[0]
            ax.annotate(f'최대: {top_label} ({top_val}%)',
                        xy=(top_val, n - 1), fontsize=9, color='#1E3A8A', fontweight='bold',
                        xytext=(top_val * 0.6, n - 0.3))
            fig.tight_layout()
            figures.append(fig)

        else:
            fig, ax = plt.subplots(figsize=(11, 5))
            bar_colors = pastel[:n]

            max_val = max(ds['values'])
            max_idx = ds['values'].index(max_val)
            highlight_colors = []
            for j in range(n):
                highlight_colors.append('#1E3A8A' if j == max_idx else bar_colors[j])

            bars = ax.bar(ds['labels'], ds['values'],
                          color=highlight_colors, width=0.6, edgecolor='white', linewidth=1.5)
            for j, (bar, val) in enumerate(zip(bars, ds['values'])):
                color = '#1E3A8A' if j == max_idx else '#374151'
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max_val * 0.02,
                        f'{val:,.0f}', ha='center', va='bottom',
                        fontsize=12, fontweight='bold', color=color)

            _style_ax(ax)
            ax.spines['bottom'].set_visible(True)
            ax.spines['bottom'].set_color('#E5E7EB')
            ax.set_title(ds['title'], fontsize=15, fontweight='bold', pad=18, color='#111827')
            unit_text = f'(단위: {ds["unit"]})'
            ax.text(0.99, 1.02, unit_text, transform=ax.transAxes,
                    ha='right', fontsize=9, color='#9CA3AF')
            fig.tight_layout()
            figures.append(fig)

    return figures


def generate_charts(datasets, output_dir="outputs"):
    os.makedirs(output_dir, exist_ok=True)
    figures = generate_charts_figures(datasets)
    chart_paths = []
    for i, fig in enumerate(figures):
        filename = f"chart_{i+1}.png"
        filepath = os.path.join(output_dir, filename)
        fig.savefig(filepath, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        chart_paths.append(filename)
    return chart_paths


def generate_report(notebook_content, report_type="경영진", title="", purpose="", audience="", length=""):
    consulting_base = """당신은 McKinsey, BCG, Deloitte 출신의 데이터 스토리텔링 전문가입니다.
업로드된 분석 결과를 보고 '코드 설명서'가 아닌 '의사결정을 위한 비즈니스 보고서'를 작성하세요.

[절대 금지]
- 코드 블록, 라이브러리명, 변수명, DataFrame 출력, Raw 데이터 나열
- 단순 수치 나열 (❌ "A 매출 120억, B 매출 80억")
- "분석 결과 ~로 나타났다" 식의 수동적 서술

[필수 원칙]
- 모든 수치는 비즈니스 맥락으로 해석 (✅ "A그룹이 전체 매출의 60%를 차지하며, 마케팅 예산 집중 시 ROI 개선 가능성이 높음")
- 핵심 수치는 **굵게** 강조, 변화는 ▲/▼ 표기
- 각 섹션 시작에 핵심 메시지를 큰따옴표로 (예: "디지털 채널이 성장을 주도하고 있다")
- 비교 데이터는 마크다운 표(table)로 정리
- 모든 인사이트는 Insight → Evidence → Recommendation → Expected Impact 구조

[스토리 흐름]
현황 파악 → 문제 발견 → 원인 분석 → 인사이트 도출 → 개선 방안 → 기대 효과

[결론 작성 규칙 - 가장 중요]
- "이 분석을 왜 했는지"에 대한 답을 제시
- 비즈니스 임팩트 중심 (매출/수익/비용/리스크 영향)
- 구체적 액션 아이템: 단기(1-3개월) / 중장기(6-12개월) 구분
- 리스크와 기회를 분리하여 경영 의사결정에 바로 활용 가능하도록 작성
"""

    title_rule = """
[제목 규칙 - 매우 중요]
보고서 첫 줄 # 제목에 반드시 '분석 내용의 주제'를 구체적으로 써야 합니다.
분석 데이터를 보고 어떤 분석인지 파악한 후, 그에 맞는 제목을 작성하세요.
예시: "# 2024년 상반기 신용평가 리스크 분석", "# 고객 이탈률 원인 분석 및 개선 전략", "# 매출 채널별 성과 비교 분석"
절대로 "데이터 분석 보고서" 같은 일반적인 제목을 쓰지 마세요.
날짜는 제목에 넣지 마세요 (별도 표기됨).
"""

    structure_rule = """
[보고서 구조 - 반드시 이 번호 체계를 사용하세요]
## 1. 분석 개요
### 1.1 분석 목표
### 1.2 분석 대상 및 데이터
## 2. 분석 내용
(분석 결과에 따라 2.1, 2.2, 2.3 등으로 세분화하세요. 각 소섹션마다 시각화에 적합한 핵심 데이터를 표로 정리하세요.)
### 2.1 (첫 번째 분석 주제)
### 2.2 (두 번째 분석 주제)
### 2.3 (세 번째 분석 주제 - 필요시)
## 3. 결론
### 3.1 비즈니스 결론 및 인사이트
### 3.2 권장 액션 아이템
"""

    prompts = {
        "분석요약": consulting_base + title_rule + structure_rule + """
이 보고서는 분석 요약 보고서입니다. 핵심만 빠르게 전달하세요.
2장 분석 내용은 데이터에서 발견된 주요 패턴, 비교, 추이를 소섹션별로 나눠 디테일하게 작성하세요.
각 소섹션에 반드시 마크다운 표를 포함하세요.

형식:
# (분석 내용에 맞는 구체적 제목)

## 1. 분석 개요
### 1.1 분석 목표
### 1.2 분석 대상 및 데이터
## 2. 분석 내용
### 2.1 (핵심 분석 1)
### 2.2 (핵심 분석 2)
### 2.3 (핵심 분석 3)
## 3. 결론
### 3.1 비즈니스 결론 및 인사이트
### 3.2 권장 액션 아이템
""",
        "기술분석": consulting_base + title_rule + structure_rule + """
이 보고서는 기술 분석 보고서입니다. 분석 방법론과 데이터 품질을 포함하되, 비즈니스 맥락을 잃지 마세요.
2장에서 기술적 세부사항을 디테일하게 작성하고, 각 소섹션에 마크다운 표를 포함하세요.

형식:
# (분석 내용에 맞는 구체적 제목)

## 1. 분석 개요
### 1.1 분석 목표 및 방법론
### 1.2 데이터 특성 및 품질
## 2. 분석 내용
### 2.1 (상세 분석 1)
### 2.2 (상세 분석 2)
### 2.3 (상세 분석 3)
## 3. 결론
### 3.1 비즈니스 임팩트
### 3.2 추가 분석 제안
""",
        "경영진": consulting_base + title_rule + structure_rule + """
이 보고서는 경영진 보고서입니다. 5분 안에 핵심을 파악할 수 있도록 작성하세요.
2장에서 현황과 문제를 디테일하게 진단하고, 각 소섹션에 근거 데이터를 표로 제시하세요.

형식:
# (분석 내용에 맞는 구체적 제목)

## 1. 분석 개요
### 1.1 분석 목표
### 1.2 분석 대상
## 2. 분석 내용
### 2.1 (현황 진단)
### 2.2 (핵심 이슈 분석)
### 2.3 (원인과 근거)
## 3. 결론
### 3.1 전략적 제언과 기대효과
### 3.2 Next Steps
"""
    }
    
    prompt = prompts[report_type]
    today = datetime.now().strftime("%Y년 %m월 %d일")
    prompt = prompt.replace("{날짜}", today)

    length_map = {
        "짧게 (1-2 페이지)": "A4 1-2장 이내로 간결하게",
        "중간 (3-5 페이지)": "A4 3-5장 분량으로 적절히 상세하게",
        "상세 (5-10 페이지)": "A4 5-10장 분량으로 매우 상세하게",
    }

    context_parts = []
    if title:
        context_parts.append(f"- 보고서 제목: {title}")
    if purpose:
        context_parts.append(f"- 분석 목적: {purpose}")
    if audience:
        if "외부" in audience:
            context_parts.append("- 대상 독자: 외부 공유용. 전문 용어를 최소화하고, 쉽고 친근한 표현을 사용하세요. 핵심 위주로 간결하게 작성하세요.")
        elif "내부" in audience:
            context_parts.append("- 대상 독자: 팀 내부용. 전문 용어 사용 가능하며, 상세한 분석 내용과 기술적 근거를 포함하세요.")
        else:
            context_parts.append(f"- 대상 독자: {audience}")
    if length and length in length_map:
        context_parts.append(f"- 분량: {length_map[length]}")

    context_str = ""
    if context_parts:
        context_str = "\n\n추가 요구사항:\n" + "\n".join(context_parts)

    max_tokens = 2000
    if "상세" in length:
        max_tokens = 4000
    elif "중간" in length:
        max_tokens = 3000

    import time
    for attempt in range(3):
        try:
            message = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=max_tokens,
                messages=[
                    {
                        "role": "user",
                        "content": f"{prompt}{context_str}\n\n---\n분석 결과:\n{notebook_content}"
                    }
                ]
            )
            return message.content[0].text
        except Exception as e:
            if "529" in str(e) or "overloaded" in str(e).lower():
                if attempt < 2:
                    time.sleep(5 * (attempt + 1))
                    continue
            raise


def save_report(report_content, report_type, chart_paths=None, output_dir="outputs"):
    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    filename = f"{output_dir}/DANA_{report_type}_{timestamp}.md"

    chart_section = ""
    if chart_paths:
        chart_section = "\n\n---\n\n## 시각화 자료\n\n"
        for cp in chart_paths:
            chart_section += f"![chart]({cp})\n\n"

    with open(filename, 'w', encoding='utf-8') as f:
        f.write(report_content)
        f.write(chart_section)

    return filename


def _clean_md(text):
    text = text.replace('**', '').replace('*', '')
    lines = []
    for line in text.split('\n'):
        s = line.strip()
        if s.startswith('#### '):
            lines.append('- ' + s.lstrip('# ').strip())
        elif s.startswith('# '):
            lines.append('# ' + s.lstrip('# ').strip())
        elif s.startswith('## '):
            lines.append('## ' + s.lstrip('# ').strip())
        elif s.startswith('### '):
            lines.append('### ' + s.lstrip('# ').strip())
        elif s == '---':
            continue
        else:
            lines.append(line)
    return '\n'.join(lines)


def export_docx(report_text, chart_figures=None):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    import io
    report_text = _clean_md(report_text)

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    style = doc.styles['Normal']
    style.font.name = '맑은 고딕'
    style.font.size = Pt(10)
    style.paragraph_format.space_after = Pt(3)
    style.paragraph_format.line_spacing = 1.5

    NAVY = RGBColor(0x00, 0x14, 0x89)
    NAVY_LIGHT = RGBColor(0x1D, 0x4E, 0xD8)
    DARK = RGBColor(0x11, 0x18, 0x27)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    RED = RGBColor(0xDC, 0x26, 0x26)
    BLUE = RGBColor(0x25, 0x63, 0xEB)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def _shd(element, color):
        pPr = element.get_or_add_pPr()
        shd = pPr.makeelement(qn('w:shd'), {
            qn('w:val'): 'clear', qn('w:color'): 'auto', qn('w:fill'): color,
        })
        pPr.append(shd)

    def _cell_shd(cell, color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = tcPr.makeelement(qn('w:shd'), {
            qn('w:val'): 'clear', qn('w:color'): 'auto', qn('w:fill'): color,
        })
        tcPr.append(shd)

    # ── 파싱 ──
    report_title_text = "DANA 분석 보고서"
    sections_data = []
    cur = None
    for line in report_text.split('\n'):
        s = line.strip()
        if s.startswith('# '):
            report_title_text = s[2:]
        elif s.startswith('## '):
            if cur:
                sections_data.append(cur)
            cur = {'title': s[3:], 'lines': []}
        elif cur is not None and s:
            cur['lines'].append(s)
    if cur:
        sections_data.append(cur)

    # ── 1. 표지 ──
    for _ in range(6):
        doc.add_paragraph('')

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('━' * 40)
    run.font.color.rgb = NAVY
    run.font.size = Pt(10)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(20)
    run = p.add_run(report_title_text)
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = NAVY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(16)
    run = p.add_run('━' * 40)
    run.font.color.rgb = NAVY
    run.font.size = Pt(10)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(30)
    today_str = datetime.now().strftime("%Y년 %m월 %d일")
    run = p.add_run(today_str)
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(8)
    run = p.add_run('DANA  |  Data Analysis NICE Automation')
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY

    # ── 2. 목차 ──
    doc.add_page_break()
    p = doc.add_paragraph()
    _shd(p.paragraph_format.element, '001489')
    p.paragraph_format.space_after = Pt(16)
    run = p.add_run('  목 차')
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    for i, sec in enumerate(sections_data):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.left_indent = Cm(1)
        if i % 2 == 0:
            _shd(p.paragraph_format.element, 'EFF6FF')
        run = p.add_run(f'{i+1}.  {sec["title"]}')
        run.font.size = Pt(11)
        run.font.color.rgb = DARK

    doc.add_paragraph('')

    # ── 3. 본문 ──
    chart_idx = 0
    for sec_i, sec in enumerate(sections_data):
        doc.add_page_break()

        p = doc.add_paragraph()
        _shd(p.paragraph_format.element, '001489')
        p.paragraph_format.space_after = Pt(12)
        run = p.add_run(f'  {sec["title"]}')
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE

        in_table = False
        current_table = None
        table_is_first_row = False

        for line in sec['lines']:
            if line.startswith('### '):
                in_table = False
                current_table = None
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(12)
                run = p.add_run(line[4:])
                run.font.size = Pt(11.5)
                run.font.bold = True
                run.font.color.rgb = NAVY

            elif line.startswith('"') and line.endswith('"'):
                in_table = False
                current_table = None
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(8)
                p.paragraph_format.space_after = Pt(8)
                p.paragraph_format.left_indent = Cm(0.3)
                p.paragraph_format.right_indent = Cm(0.3)
                _shd(p.paragraph_format.element, 'DBEAFE')
                run = p.add_run(f'  💡 {line}  ')
                run.font.size = Pt(10.5)
                run.font.bold = True
                run.font.color.rgb = NAVY

            elif line.startswith('|') and '|' in line[1:]:
                cells = [c.strip().replace('**', '') for c in line.split('|')[1:-1]]
                if all(set(c) <= set('- :') for c in cells):
                    continue
                if not in_table:
                    current_table = doc.add_table(rows=0, cols=len(cells))
                    current_table.alignment = WD_TABLE_ALIGNMENT.CENTER
                    current_table.autofit = True
                    in_table = True
                    table_is_first_row = True

                row = current_table.add_row()
                for ci, ct in enumerate(cells):
                    if ci < len(row.cells):
                        cell = row.cells[ci]
                        cell.text = ct
                        for paragraph in cell.paragraphs:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(9)
                                run.font.name = '맑은 고딕'
                        if table_is_first_row:
                            _cell_shd(cell, '001489')
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = WHITE
                                    run.font.bold = True
                        elif ci == 0:
                            _cell_shd(cell, 'EFF6FF')

                if table_is_first_row:
                    table_is_first_row = False

            elif line.startswith('- '):
                in_table = False
                current_table = None
                clean = line[2:].replace('**', '')
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.left_indent = Cm(1.5)

                run = p.add_run(clean)
                run.font.size = Pt(10)
                run.font.name = '맑은 고딕'
                if '▲' in clean or '증가' in clean:
                    run.font.color.rgb = RED
                elif '▼' in clean or '감소' in clean:
                    run.font.color.rgb = BLUE
                else:
                    run.font.color.rgb = DARK

            else:
                in_table = False
                current_table = None
                p = doc.add_paragraph()
                remaining = line
                while '**' in remaining:
                    before, _, after = remaining.partition('**')
                    if before:
                        run = p.add_run(before)
                        run.font.size = Pt(10)
                        run.font.color.rgb = DARK
                    if '**' in after:
                        bold_text, _, remaining = after.partition('**')
                        run = p.add_run(bold_text)
                        run.font.size = Pt(10)
                        run.font.bold = True
                        run.font.color.rgb = NAVY
                    else:
                        remaining = after
                if remaining:
                    clean = remaining.replace('> ', '')
                    run = p.add_run(clean)
                    run.font.size = Pt(10)
                    run.font.color.rgb = DARK

        if chart_figures and chart_idx < len(chart_figures):
            doc.add_paragraph('')
            img_buf = io.BytesIO()
            chart_figures[chart_idx].savefig(img_buf, format='png', dpi=150,
                                             bbox_inches='tight', facecolor='white')
            img_buf.seek(0)
            doc.add_picture(img_buf, width=Inches(5.8))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            chart_idx += 1

    # ── 남은 차트 ──
    if chart_figures and chart_idx < len(chart_figures):
        doc.add_page_break()
        p = doc.add_paragraph()
        _shd(p.paragraph_format.element, '001489')
        p.paragraph_format.space_after = Pt(12)
        run = p.add_run('  시각화 자료')
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE

        for fig in chart_figures[chart_idx:]:
            doc.add_paragraph('')
            img_buf = io.BytesIO()
            fig.savefig(img_buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            img_buf.seek(0)
            doc.add_picture(img_buf, width=Inches(5.8))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def export_xlsx(report_text, datasets=None):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
    from openpyxl.utils import get_column_letter
    import io
    report_text = _clean_md(report_text)

    wb = Workbook()
    NAVY = "001489"
    NAVY_LIGHT = "1D4ED8"
    BG_LIGHT = "DBEAFE"
    BG_LIGHTER = "EFF6FF"
    GRAY = "6B7280"
    DARK = "111827"
    WHITE = "FFFFFF"

    navy_fill = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
    light_fill = PatternFill(start_color=BG_LIGHT, end_color=BG_LIGHT, fill_type="solid")
    lighter_fill = PatternFill(start_color=BG_LIGHTER, end_color=BG_LIGHTER, fill_type="solid")
    white_fill = PatternFill(start_color=WHITE, end_color=WHITE, fill_type="solid")
    thin_border = Border(
        bottom=Side(style='thin', color='D1D5DB'),
        top=Side(style='thin', color='D1D5DB'),
        left=Side(style='thin', color='D1D5DB'),
        right=Side(style='thin', color='D1D5DB'),
    )

    title_font = Font(bold=True, size=16, color=WHITE, name='맑은 고딕')
    h2_font = Font(bold=True, size=12, color=NAVY, name='맑은 고딕')
    insight_font = Font(italic=True, size=10.5, color=NAVY_LIGHT, name='맑은 고딕')
    normal_font = Font(size=10, color=DARK, name='맑은 고딕')
    bullet_font = Font(size=10, color=DARK, name='맑은 고딕')
    table_header_font = Font(bold=True, size=10, color=WHITE, name='맑은 고딕')
    table_cell_font = Font(size=10, color=DARK, name='맑은 고딕')
    wrap_align = Alignment(wrap_text=True, vertical='top')

    # ── 보고서 요약 시트 ──
    ws = wb.active
    ws.title = "보고서 요약"
    ws.sheet_properties.tabColor = NAVY

    for col in range(1, 9):
        ws.column_dimensions[get_column_letter(col)].width = 16
    ws.column_dimensions['A'].width = 4
    ws.column_dimensions['B'].width = 50
    ws.column_dimensions['C'].width = 20
    ws.column_dimensions['D'].width = 20
    ws.column_dimensions['E'].width = 20
    ws.column_dimensions['F'].width = 20

    sections = []
    current = None
    for line in report_text.split('\n'):
        s = line.strip()
        if s.startswith('# '):
            report_title_text = s[2:]
        elif s.startswith('## '):
            if current:
                sections.append(current)
            current = {'title': s[3:], 'insights': [], 'bullets': [], 'table_rows': []}
        elif current:
            if s.startswith('"') and s.endswith('"'):
                current['insights'].append(s)
            elif s.startswith('|') and '|' in s[1:]:
                cells = [c.strip().replace('**', '') for c in s.split('|')[1:-1]]
                if not all(set(c) <= set('- :') for c in cells):
                    current['table_rows'].append(cells)
            elif s.startswith('- '):
                current['bullets'].append(s[2:].replace('**', ''))
            elif s and not s.startswith('---'):
                current['bullets'].append(s.replace('**', '').replace('> ', ''))
    if current:
        sections.append(current)

    row = 1
    for col in range(1, 7):
        ws.cell(row=row, column=col).fill = navy_fill
    title_cell = ws.cell(row=row, column=2, value=report_title_text if 'report_title_text' in dir() else 'DANA 분석 보고서')
    title_cell.font = title_font
    title_cell.alignment = Alignment(vertical='center')
    ws.row_dimensions[row].height = 40
    row += 2

    for sec in sections:
        for col in range(1, 7):
            ws.cell(row=row, column=col).fill = light_fill
        sec_cell = ws.cell(row=row, column=2, value=sec['title'])
        sec_cell.font = h2_font
        ws.row_dimensions[row].height = 28
        row += 1

        for ins in sec['insights']:
            ws.cell(row=row, column=2, value=ins).font = insight_font
            for col in range(1, 7):
                ws.cell(row=row, column=col).fill = lighter_fill
            row += 1

        if sec['table_rows']:
            n_cols = max(len(r) for r in sec['table_rows'])
            for i, tr in enumerate(sec['table_rows']):
                for j, val in enumerate(tr):
                    c = ws.cell(row=row, column=j+2, value=val)
                    c.border = thin_border
                    if i == 0:
                        c.font = table_header_font
                        c.fill = navy_fill
                        c.alignment = Alignment(horizontal='center', vertical='center')
                    else:
                        c.font = table_cell_font
                        c.alignment = Alignment(vertical='center')
                        if j == 0:
                            c.fill = lighter_fill
                row += 1
            row += 1

        for b in sec['bullets']:
            ws.cell(row=row, column=2, value=f"  •  {b}").font = bullet_font
            ws.cell(row=row, column=2).alignment = wrap_align
            row += 1

        row += 1

    # ── 데이터 시트들 ──
    if datasets:
        for ds in datasets:
            raw_name = ds['title']
            for ch in [':', '\\', '/', '?', '*', '[', ']']:
                raw_name = raw_name.replace(ch, ' ')
            ws_name = raw_name[:28].strip()
            ws_data = wb.create_sheet(title=ws_name)
            ws_data.sheet_properties.tabColor = NAVY_LIGHT

            ws_data.column_dimensions['A'].width = 4
            ws_data.column_dimensions['B'].width = 25
            ws_data.column_dimensions['C'].width = 18
            ws_data.column_dimensions['D'].width = 18

            for col in range(1, 5):
                ws_data.cell(row=1, column=col).fill = navy_fill
            t = ws_data.cell(row=1, column=2, value=ds['title'])
            t.font = Font(bold=True, size=13, color=WHITE, name='맑은 고딕')
            ws_data.row_dimensions[1].height = 35

            headers = ['항목', f'값 ({ds["unit"]})', '비중']
            total = sum(ds['values']) if ds['values'] else 1
            for j, h in enumerate(headers):
                c = ws_data.cell(row=3, column=j+2, value=h)
                c.font = table_header_font
                c.fill = navy_fill
                c.border = thin_border
                c.alignment = Alignment(horizontal='center')

            max_val = max(ds['values']) if ds['values'] else 0
            for k, (label, value) in enumerate(zip(ds['labels'], ds['values'])):
                r = k + 4
                c1 = ws_data.cell(row=r, column=2, value=label)
                c2 = ws_data.cell(row=r, column=3, value=value)
                c3 = ws_data.cell(row=r, column=4, value=value/total if total else 0)

                c1.font = table_cell_font
                c2.font = table_cell_font
                c3.font = table_cell_font
                c2.number_format = '#,##0.##'
                c3.number_format = '0.0%'

                for c in [c1, c2, c3]:
                    c.border = thin_border

                if value == max_val:
                    highlight = PatternFill(start_color=BG_LIGHT, end_color=BG_LIGHT, fill_type="solid")
                    for c in [c1, c2, c3]:
                        c.fill = highlight
                        c.font = Font(bold=True, size=10, color=NAVY, name='맑은 고딕')
                elif k % 2 == 0:
                    for c in [c1, c2, c3]:
                        c.fill = lighter_fill

            sum_row = len(ds['labels']) + 4
            ws_data.cell(row=sum_row, column=2, value="합계").font = Font(bold=True, size=10, color=NAVY, name='맑은 고딕')
            sc = ws_data.cell(row=sum_row, column=3, value=total)
            sc.font = Font(bold=True, size=10, color=NAVY, name='맑은 고딕')
            sc.number_format = '#,##0.##'
            ws_data.cell(row=sum_row, column=4, value=1).font = Font(bold=True, size=10, color=NAVY, name='맑은 고딕')
            ws_data.cell(row=sum_row, column=4).number_format = '0.0%'
            for col in range(2, 5):
                ws_data.cell(row=sum_row, column=col).border = Border(
                    top=Side(style='double', color=NAVY),
                    bottom=Side(style='double', color=NAVY),
                )

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def export_pptx(report_text, chart_figures=None, title="", purpose="",
                audience="", dept="", author="", datasets=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    report_text = _clean_md(report_text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    NAVY = RGBColor(0x00, 0x14, 0x89)
    NAVY_M = RGBColor(0x1D, 0x4E, 0xD8)
    DARK = RGBColor(0x11, 0x18, 0x27)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    RED = RGBColor(0xDC, 0x26, 0x26)
    BLU = RGBColor(0x25, 0x63, 0xEB)
    BG = RGBColor(0xDB, 0xEA, 0xFE)

    today_str = datetime.now().strftime("%Y.%m.%d")
    auto_title = ""
    for line in report_text.split('\n'):
        s = line.strip()
        if s.startswith('# ') and not s.startswith('## '):
            auto_title = s[2:].strip()
            auto_title = re.sub(r'\d{4}년\s*\d{1,2}월\s*\d{1,2}일', '', auto_title).strip()
            break
    dtitle = title if title else (auto_title if auto_title else "데이터 분석 보고서")

    def _bg(s, r, g, b):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(r, g, b)
    def _rect(s, l, t, w, h, c):
        from pptx.enum.shapes import MSO_SHAPE
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    def _tb(s, l, t, w, h, txt, sz=18, b=False, c=DARK, a=PP_ALIGN.LEFT):
        bx = s.shapes.add_textbox(l, t, w, h); bx.text_frame.word_wrap = True
        p = bx.text_frame.paragraphs[0]; p.text = str(txt)[:150]
        p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c
        p.font.name = '맑은 고딕'; p.alignment = a
    def _multi(s, l, t, w, h, items, sz=10):
        bx = s.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
        for i, (txt, bold, color) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt[:120]; p.font.size = Pt(sz); p.font.bold = bold
            p.font.color.rgb = color; p.font.name = '맑은 고딕'; p.space_after = Pt(3)
    def _hdr(s, txt, pg):
        _rect(s, Inches(0), Inches(0), W, Inches(1.0), NAVY)
        _rect(s, Inches(0), Inches(1.0), W, Emu(18000), NAVY_M)
        _tb(s, Inches(0.7), Inches(0.15), Inches(11), Inches(0.7), txt, sz=22, b=True, c=WHITE)
        _tb(s, W-Inches(1), H-Inches(0.4), Inches(0.6), Inches(0.3), str(pg), sz=8, c=GRAY, a=PP_ALIGN.RIGHT)
    def _tbl(sl, rows, x, y, mw=Inches(11)):
        if not rows: return Inches(0)
        nc = max(len(r) for r in rows); nr = min(len(rows), 8)
        cw = min(Inches(2.8), mw // nc)
        ts = sl.shapes.add_table(nr, nc, x, y, cw*nc, Inches(0.34)*nr)
        for ri in range(nr):
            for ci, ct in enumerate(rows[ri]):
                if ci >= nc: continue
                cell = ts.table.cell(ri, ci); cell.text = ct
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9); p.font.name = '맑은 고딕'; p.alignment = PP_ALIGN.CENTER
                    if ri == 0: p.font.bold = True; p.font.color.rgb = WHITE
                if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                elif ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEF,0xF6,0xFF)
        return Inches(0.34)*nr + Inches(0.2)
    def _chart(sl, fig, x, y, w, h):
        ib = io.BytesIO(); fig.savefig(ib, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        ib.seek(0); sl.shapes.add_picture(ib, x, y, w, h)

    # ── 파싱 ──
    mains = []; cur = None; cursub = None
    for line in report_text.split('\n'):
        s = line.strip()
        if s.startswith('# ') and not s.startswith('## '): continue
        elif s.startswith('## '):
            if cursub and cur: cur['subs'].append(cursub)
            if cur: mains.append(cur)
            cur = {'title': s[3:].strip(), 'subs': [], 'lines': []}; cursub = None
        elif s.startswith('### '):
            if cursub and cur: cur['subs'].append(cursub)
            cursub = {'title': s[4:].strip(), 'ins': [], 'bul': [], 'tbl': []}
        elif cursub and s:
            if s.startswith('"') and s.endswith('"'): cursub['ins'].append(s.strip('"'))
            elif s.startswith('|') and '|' in s[1:]:
                cells = [c.strip() for c in s.split('|')[1:-1]]
                if not all(set(c) <= set('- :') for c in cells): cursub['tbl'].append(cells)
            elif s.startswith('- '): cursub['bul'].append(s[2:])
            else: cursub['bul'].append(s)
        elif cur and s and not cursub:
            if s.startswith('"') and s.endswith('"'): cur['lines'].append(('ins', s.strip('"')))
            elif s.startswith('- '): cur['lines'].append(('bul', s[2:]))
            else: cur['lines'].append(('bul', s))
    if cursub and cur: cur['subs'].append(cursub)
    if cur: mains.append(cur)

    ov, dt, cn = [], [], []
    for sec in mains:
        t = sec['title']
        if re.match(r'^1[\.\s]', t) or '개요' in t or '목표' in t: ov.append(sec)
        elif re.match(r'^3[\.\s]', t) or '결론' in t or '제언' in t or 'Next' in t: cn.append(sec)
        else: dt.append(sec)
    if not ov and mains: ov = [mains[0]]
    if not cn and len(mains) >= 3: cn = [mains[-1]]
    if not dt: dt = [s for s in mains if s not in ov and s not in cn]

    pg = 0

    # 표지
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
    _rect(sl, Inches(0), Inches(0), W, Inches(0.12), NAVY)
    _rect(sl, Inches(0), H-Inches(0.12), W, Inches(0.12), NAVY)
    _tb(sl, Inches(0), Inches(2.0), W, Inches(0.4), "DANA Report", sz=14, c=NAVY_M, a=PP_ALIGN.CENTER)
    _rect(sl, Inches(5.2), Inches(2.5), Inches(3), Emu(14000), NAVY)
    _tb(sl, Inches(0), Inches(2.8), W, Inches(1.5), dtitle, sz=34, b=True, c=DARK, a=PP_ALIGN.CENTER)
    _tb(sl, Inches(0), Inches(4.8), W, Inches(0.4), "  |  ".join([x for x in [today_str, dept, author] if x]),
        sz=13, c=GRAY, a=PP_ALIGN.CENTER)

    # 목차
    pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
    _hdr(sl, "목차  Contents", pg)
    for i, (n, t, d) in enumerate([("01","분석 개요","분석 목표 · 데이터 개요"),
                                    ("02","분석 내용","디테일 분석 · 시각화"),
                                    ("03","결론","비즈니스 결론 · 액션 아이템")]):
        y = Inches(1.5)+Inches(i*1.5)
        _rect(sl, Inches(0.8), y, Inches(1.0), Inches(1.0), NAVY)
        _tb(sl, Inches(0.8), y+Inches(0.18), Inches(1.0), Inches(0.6), n, sz=26, b=True, c=WHITE, a=PP_ALIGN.CENTER)
        _tb(sl, Inches(2.2), y+Inches(0.1), Inches(9), Inches(0.4), t, sz=18, b=True, c=DARK)
        _tb(sl, Inches(2.2), y+Inches(0.5), Inches(9), Inches(0.3), d, sz=11, c=GRAY)

    # 분석 개요
    for sec in ov:
        pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
        _hdr(sl, f"01  {sec['title']}", pg); y = Inches(1.3)
        for sub in sec['subs']:
            if y > Inches(5.8):
                pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
                _hdr(sl, f"01  {sub['title']}", pg); y = Inches(1.3)
            _tb(sl, Inches(0.6), y, Inches(11), Inches(0.3), sub['title'], sz=13, b=True, c=NAVY); y += Inches(0.38)
            for ins in sub['ins'][:1]:
                _rect(sl, Inches(0.5), y, Inches(12), Inches(0.42), BG)
                _tb(sl, Inches(0.8), y+Inches(0.05), Inches(11.5), Inches(0.32), f"💡 {ins}", sz=10, b=True, c=NAVY)
                y += Inches(0.5)
            for b in sub['bul'][:5]:
                col = RED if '▲' in b else BLU if '▼' in b else DARK
                _tb(sl, Inches(0.8), y, Inches(11), Inches(0.28), f"• {b[:100]}", sz=9, c=col); y += Inches(0.3)

    # 분석 내용
    ci = 0
    for sec in dt:
        subs = sec['subs'] if sec['subs'] else [{'title':sec['title'],'ins':[],'bul':[l[1] for l in sec['lines']],'tbl':[]}]
        for sub in subs:
            pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
            _hdr(sl, f"02  {sub['title']}", pg)
            hc = chart_figures and ci < len(chart_figures); ht = len(sub['tbl']) > 0; y = Inches(1.2)
            if sub['ins']:
                _rect(sl, Inches(0.5), y, Inches(12), Inches(0.42), BG)
                _tb(sl, Inches(0.8), y+Inches(0.05), Inches(11.5), Inches(0.32), f"💡 {sub['ins'][0]}", sz=10, b=True, c=NAVY)
                y += Inches(0.55)
            if hc and ht:
                _tbl(sl, sub['tbl'][:6], Inches(0.5), y, mw=Inches(5.5))
                _chart(sl, chart_figures[ci], Inches(6.5), Inches(1.2), Inches(6.3), Inches(5)); ci += 1
            elif hc:
                items = [(f"• {b[:90]}", False, RED if '▲' in b else BLU if '▼' in b else DARK) for b in sub['bul'][:6]]
                if items: _multi(sl, Inches(0.6), y, Inches(5.5), Inches(4.5), items)
                _chart(sl, chart_figures[ci], Inches(6.3), Inches(1.2), Inches(6.5), Inches(5.2)); ci += 1
            elif ht:
                y += _tbl(sl, sub['tbl'], Inches(0.6), y)
                for b in sub['bul'][:5]:
                    col = RED if '▲' in b else BLU if '▼' in b else DARK
                    _tb(sl, Inches(0.8), y, Inches(11), Inches(0.28), f"• {b[:100]}", sz=9, c=col); y += Inches(0.3)
            else:
                for b in sub['bul'][:10]:
                    col = RED if '▲' in b else BLU if '▼' in b else DARK
                    _tb(sl, Inches(0.7), y, Inches(11.5), Inches(0.28), f"• {b[:100]}", sz=9, c=col); y += Inches(0.3)
                    if y > Inches(6.3): break

    if chart_figures and ci < len(chart_figures):
        for i in range(ci, len(chart_figures), 2):
            pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
            _hdr(sl, "02  시각화 자료", pg)
            for j in range(2):
                if i+j < len(chart_figures):
                    _chart(sl, chart_figures[i+j], Inches(0.4)+Inches(j*6.4), Inches(1.3), Inches(6), Inches(5.5))

    # 결론
    for sec in cn:
        pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
        _hdr(sl, f"03  {sec['title']}", pg); y = Inches(1.3)
        if purpose:
            _rect(sl, Inches(0.5), y, Inches(12), Inches(0.42), RGBColor(0xEF,0xF6,0xFF))
            _tb(sl, Inches(0.8), y+Inches(0.05), Inches(11.5), Inches(0.32),
                f"📌 분석 목적: {purpose}", sz=10, b=True, c=NAVY); y += Inches(0.55)
        for sub in sec['subs']:
            if y > Inches(5.8):
                pg += 1; sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0xFF, 0xFF, 0xFF)
                _hdr(sl, f"03  {sub['title']}", pg); y = Inches(1.3)
            _tb(sl, Inches(0.6), y, Inches(11), Inches(0.3), sub['title'], sz=13, b=True, c=NAVY); y += Inches(0.38)
            for ins in sub['ins'][:1]:
                _rect(sl, Inches(0.5), y, Inches(12), Inches(0.42), BG)
                _tb(sl, Inches(0.8), y+Inches(0.05), Inches(11.5), Inches(0.32), f"💡 {ins}", sz=10, b=True, c=NAVY)
                y += Inches(0.5)
            if sub['tbl']: y += _tbl(sl, sub['tbl'][:6], Inches(0.6), y)
            for b in sub['bul'][:6]:
                col = RED if '▲' in b else BLU if '▼' in b else DARK
                _tb(sl, Inches(0.8), y, Inches(11), Inches(0.28), f"• {b[:100]}", sz=9, c=col); y += Inches(0.3)
        if not sec['subs']:
            for typ, txt in sec['lines'][:10]:
                if typ == 'ins':
                    _rect(sl, Inches(0.5), y, Inches(12), Inches(0.42), BG)
                    _tb(sl, Inches(0.8), y+Inches(0.05), Inches(11.5), Inches(0.32), f"💡 {txt}", sz=10, b=True, c=NAVY)
                    y += Inches(0.5)
                else:
                    col = RED if '▲' in txt else BLU if '▼' in txt else DARK
                    _tb(sl, Inches(0.8), y, Inches(11), Inches(0.28), f"• {txt[:100]}", sz=9, c=col); y += Inches(0.3)

    # 감사합니다
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, 0x00, 0x14, 0x89)
    _tb(sl, Inches(0), Inches(2.0), W, Inches(1), "감사합니다", sz=48, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    _tb(sl, Inches(0), Inches(3.2), W, Inches(0.5), "Thank You", sz=18, c=RGBColor(0x60,0xA5,0xFA), a=PP_ALIGN.CENTER)
    _tb(sl, Inches(0), Inches(4.3), W, Inches(0.4), "  |  ".join([x for x in [dept, author, today_str] if x]),
        sz=13, c=RGBColor(0x93,0xC5,0xFD), a=PP_ALIGN.CENTER)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.getvalue()
